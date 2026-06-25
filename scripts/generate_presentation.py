#!/usr/bin/env python3
"""Generate AgroElevate Final BE Project Presentation (PPTX)."""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "docs" / "presentation" / "assets"
OUT_PPTX = ROOT / "AgroElevate_Final_Presentation.pptx"
OUT_NOTES = ROOT / "Speaker_Notes.md"
LIVE_URL = "https://agro-fair-chain.vercel.app"
GITHUB_URL = "https://github.com/AmoghRaidurg/agro-fair-chain"

COLLEGE = (
    "Pune Vidyarthi Griha's College of Engineering, "
    "Technology & Management, Pune – 411009"
)
DEPARTMENT = "Department of Electronics and Telecommunication Engineering"
UNIVERSITY = "Savitribai Phule Pune University"
GUIDE = "Dr. V. B. Thakare"
ACADEMIC_YEAR = "2025–2026"
TEAM = [
    ("Harshavardhan Kale", "B190070239"),
    ("Siddheya Masurkar", "B190070261"),
    ("Amogh Raidurg", "B190070288"),
    ("Mohit Sonsale", "B190070306"),
]
PROJECT_TITLE = (
    "AGROELEVATE: AI-POWERED AGRI-COMMERCE PLATFORM "
    "FOR FARMER TRANSACTIONS AND ANALYTICS"
)
LOGO_FILE = "pvg_college_logo.png"

# Agriculture-inspired dark + green palette
BG_DARK = RGBColor(0x0B, 0x1A, 0x12)
BG_CARD = RGBColor(0x14, 0x3D, 0x2A)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
GREEN_DARK = RGBColor(0x16, 0xA3, 0x4A)
TEXT = RGBColor(0xF0, 0xFD, 0xF4)
MUTED = RGBColor(0xA7, 0xF3, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF5, 0xC5, 0x42)

W = Inches(13.333)
H = Inches(7.5)


def set_slide_bg(slide, color: RGBColor = BG_DARK) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.08)) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, top, W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()


def add_title_block(slide, title: str, subtitle: str = "", y=Inches(0.35)) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), y, Inches(12), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(6)


def add_bullets(slide, items: list[str], x=Inches(0.7), y=Inches(1.5), w=Inches(5.5), size=18) -> None:
    box = slide.shapes.add_textbox(x, y, w, Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def add_card(slide, x, y, w, h, title: str, body: str = "") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = GREEN_DARK
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT
        p2.space_before = Pt(6)


def add_image_if_exists(slide, filename: str, x, y, w, h=None) -> bool:
    path = ASSETS / filename
    if path.exists():
        if h is None:
            slide.shapes.add_picture(str(path), x, y, width=w)
        else:
            slide.shapes.add_picture(str(path), x, y, width=w, height=h)
        return True
    return False


def generate_qr_code() -> None:
    try:
        import qrcode

        ASSETS.mkdir(parents=True, exist_ok=True)
        img = qrcode.make(LIVE_URL)
        img.save(ASSETS / "qr_live.png")
    except ImportError:
        pass


SPEAKER_NOTES: list[str] = [
    """SLIDE 1 — TITLE
Good morning/afternoon respected examiners, faculty, and guests.

I am [Student Name], presenting our Final Year BE project: AgroElevate — an AI-powered agricultural marketplace with automated royalty distribution.

AgroElevate connects farmers, traders, industrialists, and customers on a single platform where downstream resale value is shared fairly with the original producer through a server-side royalty engine.

This presentation covers architecture, implementation, live deployment, and verified results from our production system at agro-fair-chain.vercel.app.""",

    """SLIDE 2 — PROBLEM STATEMENT
Indian agriculture employs over half our workforce, yet farmers capture the smallest share of value in multi-hop supply chains.

Middlemen and processors absorb margins while the farmer receives only a one-time farm-gate payment. There is no transparency, no traceability, and no automated mechanism for revenue sharing on resale.

Existing AgriTech platforms focus on listings and logistics — not on encoding fair economics into the transaction layer itself. AgroElevate was designed to close this gap.""",

    """SLIDE 3 — OUR SOLUTION
AgroElevate is a multi-role marketplace serving five personas: Farmer, Trader, Industrialist, Customer, and Admin.

Each role has dedicated dashboards, wallet operations, and marketplace permissions. The platform pillars are: Marketplace commerce, Razorpay wallet top-ups, the Option B royalty engine at 12.5%, AI intelligence dashboards, and role-aware analytics.

All business logic — especially royalty calculation — runs server-side in PostgreSQL RPC functions, never on the client.""",

    """SLIDE 4 — SYSTEM ARCHITECTURE
The architecture follows a thin-client, authoritative-backend pattern.

React web and Android mobile clients communicate with Supabase for authentication, PostgreSQL data, Row Level Security, RPC functions, Edge Functions, and file storage.

The FastAPI AI microservice on Render provides predictions, recommendations, and the Copilot. Razorpay handles payment intents and webhook settlement. Vercel hosts the production web frontend.

This separation ensures security, scalability, and a single source of truth for commerce.""",

    """SLIDE 5 — TECH STACK
We chose industry-standard, production-grade technologies.

Frontend: React 18, TypeScript, Vite, Tailwind, and shadcn/ui for a modern responsive UI.

Backend: Supabase with PostgreSQL 15, RLS policies, and SECURITY DEFINER RPCs.

AI: FastAPI with Python, scikit-learn, and pandas, containerized for Render deployment.

Payments: Razorpay with Edge Function webhooks. DevOps: GitHub, Vercel, Render, and Docker.""",

    """SLIDE 6 — CORE FEATURES
AgroElevate delivers twelve integrated capabilities: Marketplace listings and checkout, wallet with Razorpay top-up, automated royalty engine, order lifecycle management, role-based analytics, AI intelligence layer, grounded Copilot chat, Supabase authentication, admin dashboard, real-time notifications, Android client architecture, and full Razorpay payment integration.

Each feature is role-aware and backed by server-side validation.""",

    """SLIDE 7 — ROYALTY ENGINE (CORE INNOVATION)
This is our distinguishing innovation — Option B royalty remittance.

When a Trader relists produce with embedded ownership metadata and an Industrialist purchases at resale, 12.5% of the transaction value is automatically credited to the original Farmer's wallet as royalty_income.

The calculation happens inside the checkout_order RPC function — never on the client. We verified this with 26 automated end-to-end tests, including mathematical validation of ₹43.75 royalty on a 5 kg × ₹70 resale scenario.

This encodes fair agricultural economics directly in database logic.""",

    """SLIDE 8 — AI INTELLIGENCE
The AgroElevate Intelligence API is a FastAPI microservice delivering crop recommendations, market price predictions, income forecasts in three scenarios, district-level analytics, and demand intelligence.

The AI Copilot provides grounded, context-aware responses using live user and marketplace data. The web client includes graceful offline fallback when the AI service is unavailable.

Dashboard screenshots shown here are from our live Intelligence hub at /intelligence.""",

    """SLIDE 9 — DATABASE & BACKEND
Supabase PostgreSQL is the authoritative data layer. Core tables include profiles, users, products, orders, order_items, wallet_history, payment_intents, and payment_receipts.

Row Level Security ensures users access only their own data. Critical RPCs — checkout_order, get_wallet_balance, transfer_funds — run as SECURITY DEFINER functions.

Edge Functions handle Razorpay order creation and webhook settlement, maintaining payment audit trails.""",

    """SLIDE 10 — ANDROID APPLICATION
The Android client follows the same thin-client architecture as the web app. Built with Kotlin, Jetpack Compose, and the Supabase SDK, it consumes identical RPCs and authentication flows.

Features include marketplace browsing, wallet operations, order tracking, and Razorpay Android SDK integration. The client polls payment_intents after checkout and never duplicates royalty or wallet logic locally.

Integration contracts are documented in our ANDROID_RAZORPAY_INTEGRATION.md specification.""",

    """SLIDE 11 — LIVE DEMONSTRATION
These screenshots are captured from our production deployment at agro-fair-chain.vercel.app.

The flow demonstrates: landing page, authentication, role dashboard, marketplace browsing, wallet with transaction history including royalty entries, and the AI Intelligence hub.

We invite you to explore the live system during the Q&A session. All commerce flows have been verified against the production Supabase backend.""",

    """SLIDE 12 — DEPLOYMENT
Our deployment pipeline is multi-cloud and CI/CD ready.

Source code is on GitHub. Vercel auto-deploys the React frontend on every push to main. The AI service runs as a Docker container on Render. Supabase Cloud hosts Auth, PostgreSQL, Edge Functions, and Storage.

Razorpay Edge Functions process payments in production. Release tag v1.0.0 marks our public release candidate.""",

    """SLIDE 13 — RESULTS
Key verified achievements: 26 out of 26 commerce end-to-end tests passing, 12.5% royalty rate mathematically verified, five user roles fully supported, release candidate v1.0.0 tagged and deployed, Lighthouse SEO score of 100, and zero client-side royalty calculation — all royalty logic is server-authoritative.

Overall project readiness score: 87 out of 100 across web, Android specification, AI, Razorpay, and the royalty engine.""",

    """SLIDE 14 — FUTURE SCOPE
Future enhancements include IoT soil and weather sensors, drone crop monitoring, satellite NDVI integration, e-NAM mandi price feeds, blockchain traceability, multilingual AI Copilot, export marketplace, FPO bulk listing, supply chain analytics, and an ML model retraining pipeline.

These extensions build on the fair-chain foundation already encoded in AgroElevate's backend.""",

    """SLIDE 15 — THANK YOU
Thank you for your attention. We welcome your questions and a live demonstration of AgroElevate.

Scan the QR code or visit agro-fair-chain.vercel.app. Source code is available at github.com/AmoghRaidurg/agro-fair-chain under release tag v1.0.0.

We are happy to demonstrate the royalty engine, wallet flow, and AI intelligence live. Thank you.""",
]


def apply_speaker_notes(prs: Presentation) -> None:
    for i, slide in enumerate(prs.slides):
        if i < len(SPEAKER_NOTES):
            slide.notes_slide.notes_text_frame.text = SPEAKER_NOTES[i]


def write_speaker_notes_md() -> None:
    lines = ["# AgroElevate — Speaker Notes\n", f"**Live demo:** {LIVE_URL}  \n**GitHub:** {GITHUB_URL}\n"]
    for i, note in enumerate(SPEAKER_NOTES, 1):
        title = note.split("\n", 1)[0]
        body = note.split("\n", 1)[1] if "\n" in note else ""
        lines.append(f"## Slide {i}: {title.replace('SLIDE ' + str(i) + ' — ', '')}\n")
        lines.append(body.strip() + "\n")
        lines.append("---\n")
    OUT_NOTES.write_text("\n".join(lines), encoding="utf-8")


def add_animations(pptx_path: Path) -> bool:
    """Add tasteful fade-in animations via PowerPoint COM."""
    try:
        import comtypes.client  # type: ignore

        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        # msoAnimTriggerAfterPrevious=3, msoAnimEffectFade=10
        for slide in presentation.Slides:
            first = True
            for shape in slide.Shapes:
                try:
                    trigger = 2 if first else 3  # OnPageClick / AfterPrevious
                    slide.TimeLine.MainSequence.AddEffect(
                        Shape=shape, effectId=10, trigger=trigger
                    )
                    first = False
                except Exception:
                    pass
        presentation.Save()
        presentation.Close()
        powerpoint.Quit()
        return True
    except Exception as e:
        print(f"Animation pass skipped: {e}")
        return False


def add_stat_card(slide, x, y, value: str, label: str) -> None:
    add_card(slide, x, y, Inches(2.6), Inches(1.35), value, label)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = int(W)
    prs.slide_height = int(H)
    blank = prs.slide_layouts[6]

    # --- SLIDE 1: Title ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    # Hero gradient block
    hero = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.08), W, Inches(4.2))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(0x0F, 0x2D, 0x1A)
    hero.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(7.5), Inches(2.2))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "AgroElevate"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = PROJECT_TITLE
    p2.font.size = Pt(16)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(12)
    meta = s.shapes.add_textbox(Inches(0.8), Inches(4.35), Inches(8), Inches(2.8))
    mtf = meta.text_frame
    mtf.word_wrap = True
    team_line_1 = f"{TEAM[0][0]} ({TEAM[0][1]})  |  {TEAM[1][0]} ({TEAM[1][1]})"
    team_line_2 = f"{TEAM[2][0]} ({TEAM[2][1]})  |  {TEAM[3][0]} ({TEAM[3][1]})"
    lines = [
        "Bachelor of Engineering — Final Year Project",
        COLLEGE,
        DEPARTMENT,
        UNIVERSITY,
        "Submitted By:",
        team_line_1,
        team_line_2,
        f"Guide: {GUIDE}  |  Academic Year: {ACADEMIC_YEAR}",
    ]
    for i, line in enumerate(lines):
        mp = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        mp.text = line
        mp.font.size = Pt(12)
        mp.font.color.rgb = MUTED
        mp.space_after = Pt(3)
    add_image_if_exists(s, LOGO_FILE, Inches(0.55), Inches(0.18), Inches(1.05))
    add_image_if_exists(s, "landing.png", Inches(8.2), Inches(0.5), Inches(4.8), Inches(3.0))

    # --- SLIDE 2: Problem ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Problem Statement", "Indian agriculture loses value across opaque supply chains")
    add_bullets(s, [
        "Farmers receive only farm-gate price — no share in downstream resale",
        "Middlemen capture arbitrage without transparent royalty",
        "Low profit visibility and fragmented digital infrastructure",
        "Lack of trust, traceability, and fair payment rails",
        "Existing AgriTech focuses on listings — not automated revenue sharing",
    ], y=Inches(1.35))
    add_card(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.2), "Supply Chain Gap",
             "Farmer → Trader → Industrialist → Retail\nValue concentrates downstream")
    add_card(s, Inches(7.0), Inches(4.0), Inches(5.5), Inches(2.0), "Impact",
             "< 15% farmer share in many commodity chains\nNo programmatic royalty on resale")

    # --- SLIDE 3: Solution ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Our Solution — AgroElevate", "Fair-chain commerce + server-side royalty + AI intelligence")
    roles = [
        ("Farmer", "List produce, earn royalties"),
        ("Trader", "Procure & relist inventory"),
        ("Industrialist", "Bulk procurement + manufacturing"),
        ("Customer", "Marketplace checkout"),
        ("Admin", "Users & payments oversight"),
    ]
    for i, (role, desc) in enumerate(roles):
        col, row = i % 3, i // 3
        add_card(s, Inches(0.6) + col * Inches(4.1), Inches(1.5) + row * Inches(1.55),
                 Inches(3.8), Inches(1.35), role, desc)
    add_card(s, Inches(0.6), Inches(4.8), Inches(12.0), Inches(1.6), "Platform Pillars",
             "Marketplace  •  Wallet (Razorpay)  •  Royalty Engine (12.5%)  •  AI Copilot  •  Multi-role Analytics")

    # --- SLIDE 4: Architecture ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "System Architecture", "Thin clients — authoritative backend")
    add_image_if_exists(s, "architecture.png", Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))

    # --- SLIDE 5: Tech Stack ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Technology Stack", "Production-grade, industry-standard tools")
    stack = [
        ("Frontend", "React 18 · TypeScript · Vite · Tailwind · shadcn/ui"),
        ("Mobile", "Android · Kotlin · Jetpack Compose · Supabase SDK"),
        ("Backend", "Supabase · PostgreSQL 15 · RLS · RPC · Edge Functions"),
        ("AI", "FastAPI · Python · scikit-learn · pandas · Render"),
        ("Payments", "Razorpay Test/Live · Webhook settlement"),
        ("DevOps", "GitHub · Vercel · Render · Docker"),
    ]
    for i, (cat, tech) in enumerate(stack):
        row, col = divmod(i, 2)
        add_card(s, Inches(0.6) + col * Inches(6.3), Inches(1.4) + row * Inches(1.55),
                 Inches(6.0), Inches(1.35), cat, tech)

    # --- SLIDE 6: Core Features ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Core Features", "End-to-end AgriTech commerce platform")
    features = [
        "Marketplace", "Wallet", "Royalty Engine", "Orders",
        "Analytics", "AI Intelligence", "Copilot", "Authentication",
        "Admin Dashboard", "Notifications", "Android App", "Razorpay",
    ]
    for i, feat in enumerate(features):
        col, row = divmod(i, 4)
        add_card(s, Inches(0.5) + col * Inches(3.15), Inches(1.45) + row * Inches(1.45),
                 Inches(2.9), Inches(1.2), feat, "")

    # --- SLIDE 7: Royalty Engine ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Royalty Engine — Core Innovation", "Server-side Option B · 12.5% verified (26/26 E2E)")
    add_image_if_exists(s, "royalty.png", Inches(0.4), Inches(1.2), Inches(7.8), Inches(5.5))
    add_bullets(s, [
        "Royalty ONLY on qualifying resale (trader → industrialist)",
        "Calculated in checkout_order RPC — never on client",
        "wallet_history type: royalty_income",
        "Verified: ₹43.75 on 5×₹70 resale scenario",
        "Ownership metadata in products.description JSON",
    ], x=Inches(8.4), y=Inches(1.5), w=Inches(4.5), size=16)

    # --- SLIDE 8: AI Intelligence ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "AI Intelligence Layer", "FastAPI microservice — role-aware dashboards")
    add_image_if_exists(s, "ai_pipeline.png", Inches(0.4), Inches(1.2), Inches(6.5), Inches(3.5))
    add_image_if_exists(s, "intelligence.png", Inches(7.2), Inches(1.2), Inches(5.6), Inches(3.5))
    add_bullets(s, [
        "Crop recommendations & market predictions",
        "Income forecasts (3 scenarios)",
        "District analytics & demand intelligence",
        "AI Copilot with grounded context",
        "Graceful offline fallback on web/mobile",
    ], y=Inches(4.9), w=Inches(12))

    # --- SLIDE 9: Database & Backend ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Database & Backend", "Supabase PostgreSQL with RLS + SECURITY DEFINER RPCs")
    add_image_if_exists(s, "er_diagram.png", Inches(0.4), Inches(1.2), Inches(7.5), Inches(5.5))
    add_bullets(s, [
        "profiles · users · products · orders",
        "order_items · wallet_history",
        "payment_intents · payment_receipts",
        "RPC: checkout_order, get_wallet_balance",
        "Edge: razorpay-create-order, webhook",
    ], x=Inches(8.2), y=Inches(1.5), w=Inches(4.5))

    # --- SLIDE 10: Android ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Android Application", "Thin client — same Supabase backend as web")
    add_image_if_exists(s, "android_nav.png", Inches(0.4), Inches(1.2), Inches(6.0), Inches(4.5))
    add_bullets(s, [
        "Kotlin + Jetpack Compose + Supabase SDK",
        "Auth · Marketplace · Wallet · Orders",
        "Razorpay Android Standard SDK",
        "Poll payment_intents after checkout",
        "Never duplicate royalty or wallet logic",
    ], x=Inches(6.8), y=Inches(1.5), w=Inches(5.8))
    add_image_if_exists(s, "marketplace.png", Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.8))

    # --- SLIDE 11: Live Demo ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Live Demonstration", "Production deployment — agro-fair-chain.vercel.app")
    shots = [
        ("landing.png", "Landing"),
        ("login.png", "Login"),
        ("dashboard.png", "Dashboard"),
        ("marketplace.png", "Marketplace"),
        ("wallet.png", "Wallet"),
        ("intelligence.png", "AI Hub"),
    ]
    for i, (file, label) in enumerate(shots):
        col, row = divmod(i, 3)
        x = Inches(0.45) + col * Inches(4.25)
        y = Inches(1.35) + row * Inches(2.95)
        if add_image_if_exists(s, file, x, y, Inches(4.0), Inches(2.35)):
            lb = s.shapes.add_textbox(x, y + Inches(2.4), Inches(4.0), Inches(0.35))
            lp = lb.text_frame.paragraphs[0]
            lp.text = label
            lp.font.size = Pt(12)
            lp.font.color.rgb = GREEN
            lp.alignment = PP_ALIGN.CENTER

    # --- SLIDE 12: Deployment ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Deployment Architecture", "CI/CD ready · multi-cloud production")
    add_image_if_exists(s, "deployment.png", Inches(0.4), Inches(1.2), Inches(7.5), Inches(4.5))
    add_bullets(s, [
        "Web: Vercel — agro-fair-chain.vercel.app",
        "AI: Render — FastAPI Docker container",
        "DB: Supabase Cloud (Auth + Postgres)",
        "Payments: Razorpay + Edge Functions",
        "Source: GitHub — AmoghRaidurg/agro-fair-chain",
    ], x=Inches(8.2), y=Inches(1.5), w=Inches(4.5))

    # --- SLIDE 13: Results ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Results & Achievements", "Verified against live production backend")
    stats = [
        ("26/26", "Commerce E2E tests"),
        ("12.5%", "Royalty rate verified"),
        ("5", "User roles supported"),
        ("v1.0.0-rc", "Release candidate"),
        ("100", "Lighthouse SEO score"),
        ("0", "Client-side royalty calc"),
    ]
    for i, (val, lbl) in enumerate(stats):
        col, row = divmod(i, 3)
        add_stat_card(s, Inches(0.6) + col * Inches(4.1), Inches(1.5) + row * Inches(1.65), val, lbl)
    add_card(s, Inches(0.6), Inches(5.0), Inches(12.0), Inches(1.8), "Project Readiness Score: 87/100",
             "Full-stack: Web + Android spec + AI + Razorpay + Option B royalty engine")

    # --- SLIDE 14: Future Scope ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_accent_bar(s)
    add_title_block(s, "Future Scope", "Extending fair-chain agriculture at scale")
    future = [
        "IoT soil & weather sensors", "Drone crop monitoring",
        "Satellite NDVI integration", "e-NAM mandi price feeds",
        "Blockchain traceability", "Multilingual AI Copilot",
        "Export marketplace", "FPO bulk listing",
        "Supply chain analytics", "ML model retraining pipeline",
    ]
    for i, item in enumerate(future):
        col, row = divmod(i, 2)
        add_card(s, Inches(0.6) + col * Inches(6.3), Inches(1.4) + row * Inches(1.15),
                 Inches(6.0), Inches(0.95), item, "")

    # --- SLIDE 15: Thank You ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, RGBColor(0x0F, 0x2D, 0x1A))
    add_accent_bar(s)
    t = s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    contact = s.shapes.add_textbox(Inches(0.8), Inches(2.85), Inches(9.2), Inches(3.8))
    ctf = contact.text_frame
    ctf.word_wrap = True
    team_names = ", ".join(name for name, _ in TEAM)
    clines = [
        "Questions & Live Demo Welcome",
        "",
        COLLEGE,
        DEPARTMENT,
        f"Guide: {GUIDE}  |  Academic Year: {ACADEMIC_YEAR}",
        "",
        f"Team: {team_names}",
        "",
        f"Live: {LIVE_URL}",
        f"GitHub: {GITHUB_URL}",
        "",
        "Scan QR code for live demo",
    ]
    for i, line in enumerate(clines):
        cp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        cp.text = line
        cp.font.size = Pt(14 if i > 0 else 20)
        cp.font.color.rgb = MUTED if i > 0 else TEXT
        cp.alignment = PP_ALIGN.CENTER
    add_image_if_exists(s, LOGO_FILE, Inches(0.55), Inches(0.18), Inches(1.05))
    if not add_image_if_exists(s, "qr_live.png", Inches(10.5), Inches(4.5), Inches(1.8), Inches(1.8)):
        qr = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(4.8), Inches(1.8), Inches(1.8))
        qr.fill.solid()
        qr.fill.fore_color.rgb = WHITE
        qr.line.color.rgb = GREEN
        qtf = qr.text_frame
        qtf.text = "QR"
        qtf.paragraphs[0].alignment = PP_ALIGN.CENTER
        qtf.paragraphs[0].font.size = Pt(14)
        qtf.paragraphs[0].font.color.rgb = BG_DARK

    return prs


def export_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    try:
        import comtypes.client  # type: ignore

        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        presentation.SaveAs(str(pdf_path.resolve()), 32)  # ppSaveAsPDF
        presentation.Close()
        powerpoint.Quit()
        return pdf_path.exists()
    except Exception as e:
        print(f"PDF export via PowerPoint COM failed: {e}")
        return False


def main() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    generate_qr_code()
    prs = build_presentation()
    apply_speaker_notes(prs)
    prs.save(str(OUT_PPTX))
    print(f"Created {OUT_PPTX}")

    if add_animations(OUT_PPTX):
        print("Applied fade animations")

    pdf_path = ROOT / "AgroElevate_Final_Presentation.pdf"
    if export_pdf(OUT_PPTX, pdf_path):
        print(f"Created {pdf_path}")
    else:
        print("PDF not created — open PPTX in PowerPoint and Export as PDF if needed")


if __name__ == "__main__":
    main()
